import pandas as pd
import os
from glob import glob
import numpy as np
import Plot_Generator as pg
import tkinter as tk
from tkinterdnd2 import TkinterDnD, DND_FILES
from tkinter import filedialog
import os
from pptx import Presentation
from pptx.util import Inches



class Data_Analyzer:

    def __init__(self):
        pass

    def Compile_Results(self, data_folder):

        sample_name = os.path.basename(data_folder)
        print(f"Beginning data extraction of {sample_name}")

        if not os.path.exists(data_folder + r"\\" + "Compiled Data"):
            os.makedirs(data_folder + r"\\" + "Compiled Data")
            print(f"Created folder: Compiled Data\n")
        else:
            print("Summary folder already exists, overwriting files\n")

        save_folder = data_folder + r"\\" + "Compiled Data"

        csv_files = glob(os.path.join(data_folder, "*.csv"))

        cv_df = pd.DataFrame()
        ca_df = pd.DataFrame()
        ocp_df = pd.DataFrame()
        eis_df = pd.DataFrame()

        for file_path in csv_files:
            file_name = os.path.splitext(os.path.basename(file_path))[0]

            if "-CV_" in file_name:
                data = np.loadtxt(file_path, delimiter=',')
                file_parameters = file_name.split('-CV')[1].split("_")[1:]
                meas_name = file_parameters[0] + "_" + file_parameters[1] + "_" + file_parameters[2]

                cv_df["Potential (V)"] = data[:, 0]

                for i in range(1, data.shape[1], 2):
                    cv_df[f"{meas_name}_Scan{int((i + 1) / 2)}"] = data[:, i]

                cv_df = cv_df.copy()

            elif "OCP" in file_name:

                data = np.loadtxt(file_path, delimiter=',')
                file_parameters = file_name.split('-OCP')[1].split("_")[1:]
                meas_name = file_parameters[0] + "_" + file_parameters[1] + "_" + file_parameters[2]

                ocp_df["Time (s)"] = data[:, 0]
                ocp_df[f"{meas_name}"] = data[:, 1]
                ocp_df = ocp_df.copy()

            elif "EIS" in file_name:
                data = np.loadtxt(file_path, delimiter=',')
                file_parameters = file_name.split('-EIS')[1].split("_")[1:]
                meas_name = file_parameters[0] + "_" + file_parameters[1] + "_" + file_parameters[2]

                eis_df["Frequency (hZ)"] = data[:, 2]
                eis_df[f"Zim_{meas_name}"] = data[:, 0]
                eis_df[f"Zre_{meas_name}"] = data[:, 1]
                eis_df[f"|Z|_{meas_name}"] = np.sqrt(np.power(data[:, 0], 2) + np.power(data[:, 1], 2))

                eis_df = eis_df.copy()

            elif "CA" in file_name:

                data = np.loadtxt(file_path, delimiter=',')
                file_parameters = file_name.split('-CA')[1].split("_")[1:]
                meas_name = file_parameters[0] + "_" + file_parameters[1] + "_" + file_parameters[2]

                ca_df["Time (s)"] = data[:, 0]
                ca_df[f"{meas_name}"] = data[:, 1]
                ca_df = ca_df.copy()

        ca_df.to_csv(save_folder + "\\" + "CA_Compiled.csv")
        print("Successfully saved CA data")

        cv_df.to_csv(save_folder + "\\" + "CV_Compiled.csv")
        print("Successfully saved CV data")

        eis_df.to_csv(save_folder + "\\" + "EIS_Compiled.csv")
        print("Successfully saved EIS data")

        ocp_df.to_csv(save_folder + "\\" + "OCP_Compiled.csv")
        print("Successfully saved OCP data")

        return save_folder

    def Generate_Summary_Slide(self, data_folder):

        image_folder = data_folder + r"\\" + "Summary Plots"

        if not os.path.exists(data_folder + r"\\" + "Summary Plots"):
            os.makedirs(data_folder + r"\\" + "Summary Plots")
            print(f"Created folder: Summary Plots\n")
        else:
            print("Image folder already exists, overwriting files\n")

        summary_folder = self.Compile_Results(data_folder)

        if os.path.getctime(data_folder) >= 1750557522.0:
            new_version = True
        else:
            new_version = False

        cv_images = self.Analyze_CV_Summary(summary_path=summary_folder + r"\\" + "CV_Compiled.csv", image_path=image_folder)
        ca_images = self.Analyze_CA_Summary(summary_path=summary_folder + r"\\" + "CA_Compiled.csv", new_version=new_version, image_path=image_folder)
        ocp_images = self.Analyze_OCP_Summary(summary_path=summary_folder + r"\\" + "OCP_Compiled.csv", image_path=image_folder)

        sample_name = os.path.basename(data_folder)
        template_path = r"C:\Users\Daniel\Box\Research\Software\PEC_Driver\PEC_Driver\Powerpoint_Template\Template.pptx"
        prs = Presentation(template_path)
        slide_layout = prs.slide_layouts[1]  # Use an appropriate layout index

        slide1 = prs.slides.add_slide(slide_layout)
        title1 = slide1.shapes.title
        title1.text = f"{sample_name} - CV: Wavelength"
        slide1.shapes.add_picture(cv_images[0], Inches(0.67), Inches(1.29), height=Inches(6))


        slide2 = prs.slides.add_slide(slide_layout)
        title2 = slide2.shapes.title
        title2.text = f"{sample_name} - CV: Pixels"
        slide2.shapes.add_picture(cv_images[1], Inches(0), Inches(1.5), height=Inches(5))
        #slide2.shapes.add_picture(cv_images[2], Inches(6.67), Inches(1.5), height=Inches(5))

        slide3 = prs.slides.add_slide(slide_layout)
        title3 = slide3.shapes.title
        title3.text = f"{sample_name} - CA: Summary"
        slide3.shapes.add_picture(ca_images[2], Inches(1.86), Inches(1.25), height=Inches(6.25))

        slide4 = prs.slides.add_slide(slide_layout)
        title4 = slide4.shapes.title
        title4.text = f"{sample_name} - CA: Raw and Average Traces"
        slide4.shapes.add_picture(ca_images[0], Inches(0), Inches(1.5), height=Inches(5))
        slide4.shapes.add_picture(ca_images[1], Inches(6.67), Inches(1.5), height=Inches(5))


        slide5 = prs.slides.add_slide(slide_layout)
        title5 = slide5.shapes.title
        title5.text = f"{sample_name} - OCP: Wavelength"
        slide5.shapes.add_picture(ocp_images[0], Inches(.67), Inches(1.29), height=Inches(6))


        # Save new presentation
        prs.save(summary_folder + r'\Device Summary.pptx')
    def Analyze_CA_Summary(self, summary_path, new_version=False, image_path=None):
        ca_df = pd.read_csv(summary_path, index_col=0)

        time_df = ca_df['Time (s)']
        current_df = ca_df.drop(columns=['Time (s)'])

        if not new_version:
            p1_columns = current_df.drop(
                columns=current_df.columns[[i for i, v in enumerate(current_df.columns) if "P1" not in v]])
            p1_columns = p1_columns.iloc[41:841, :].reset_index(drop=True)
            p234_columns = current_df.drop(
                columns=current_df.columns[[i for i, v in enumerate(current_df.columns) if "P1" in v]])
            p234_columns = p234_columns.iloc[50:850, :].reset_index(drop=True)
            adj_current_df = pd.concat([p1_columns, p234_columns], axis=1, ignore_index=False)
        else:
            adj_current_df = current_df.iloc[50:850, :].reset_index(drop=True)

        time_df = time_df.iloc[0:800].reset_index(drop=True)

        cur_sums, pulse_traces = self.Calculate_CA_Average(adj_current_df)

        wavelength_array = list(dict.fromkeys([col.split('nm')[0][-3:] for col in adj_current_df.columns]))
        avg_traces = pd.DataFrame()
        compiled_cur_sums = pd.DataFrame()

        for i in range(0, np.shape(pulse_traces)[0]):
            for wavelength in wavelength_array:
                temp_df = pulse_traces[i].drop(
                    columns=pulse_traces[i].columns[
                        [i for i, v in enumerate(pulse_traces[i].columns) if f"{wavelength}" not in v]])
                temp_df = temp_df.mean(axis=1)
                avg_traces[f"{wavelength}nm_P{i+1}_Power-2600"] = temp_df
                compiled_cur_sums[f"{wavelength}nm_P{i + 1}_Power-2600"] = cur_sums[i].drop(
                    columns=cur_sums[i].columns[
                        [i for i, v in enumerate(cur_sums[i].columns) if f"{wavelength}" not in v]]).values.flatten()

        if image_path is not None:
            image_1 = image_path + r"//" + 'Raw_CA_Full_Overlay.png'
            image_2 = image_path + r"//" + 'Average_CA_Overlay.png'
            image_3 = image_path + r"//" + 'Average_CA_Bar_Graph.png'
        else:
            image_1 = None
            image_2 = None
            image_3 = None

        pg.CA_Overlay_Tiled(time_df, adj_current_df, save_path=image_1, color_by='Wavelength', tile_by="Pixel",
                            target_pixels=[1, 2, 3, 4])
        pg.CA_Overlay_Tiled(time_df.iloc[0:50], avg_traces, save_path=image_2, color_by='Wavelength', tile_by="Pixel",
                            target_pixels=[1, 2, 3, 4])
        final_summary = pg.CA_Bar_Graph(compiled_cur_sums, save_path=image_3)

        final_summary.to_csv(os.path.dirname(summary_path) + "\\" + "CA_Averages.csv")

        return [image_1, image_2, image_3]
    def Analyze_CV_Summary(self, summary_path, image_path=None):

        # Section for On-Off comparison, only in reverse direction

        cv_df = pd.read_csv(summary_path, index_col=0)
        reverse_start_index = cv_df['Potential (V)'].idxmax()
        reverse_end_index = cv_df['Potential (V)'].idxmin()

        cv_reverse_df = cv_df.iloc[reverse_start_index:reverse_end_index].drop(columns="Potential (V)")
        potential_reverse_df = cv_df['Potential (V)'].iloc[reverse_start_index:reverse_end_index]

        if image_path is not None:
            image_1 = image_path + r"//" + 'Raw_CV-rev_Tiled-Wavelength.png'
            image_2 = image_path + r"//" + 'Raw_CV-rev_Tiled-Pixel.png'
        else:
            image_1 = None
            image_2 = None

        pg.CV_Overlay_Tiled(potential_reverse_df, cv_reverse_df, save_path=image_1)
        pg.CV_Overlay_Tiled(potential_reverse_df, cv_reverse_df, save_path=image_2, tile_by='Pixel',
                            color_by='Wavelength', dark_repeat=False)
        return [image_1, image_2]

    def Analyze_CVS_Summary(self, summary_path):
        pass
    def Analyze_OCP_Summary(self, summary_path, image_path=None):
        ocp_df = pd.read_csv(summary_path, index_col=0)

        time_df = ocp_df['Time (s)']
        pot_df = ocp_df.drop(columns=['Time (s)'])

        if image_path is not None:
            image_1 = image_path + r"//" + 'Raw_OCP_Tiled-Pixel.png'
        else:
            image_1 = None


        pg.OCP_Overlay_Tiled(time_df, pot_df, save_path=image_1)

        return [image_1]

    def Analyze_EIS_Summary(self, summary_path):
        pass

    def Generate_Comparison_Slide(self, target_folders, ca=True, cv=True, ocp=True, eis=True):
        template_path = r"C:\Users\Daniel\Box\Research\Software\PEC_Driver\PEC_Driver\Powerpoint_Template\Template.pptx"
        master = Presentation(template_path)
        save_path = os.path.dirname(target_folders[0]) + r"\\Device Comparison Slide.pptx"

        for folder in target_folders:
            path = folder + r"\Compiled Data\Device Summary.pptx"
            current = Presentation(path)
            for slide in current.slides:
                # Create a blank slide with same layout as source
                layout = master.slide_layouts[0]  # use layout 0 or customize
                new_slide = master.slides.add_slide(layout)

                # Copy shapes from old slide to new slide
                for shape in slide.shapes:
                    # Skip group shapes or unsupported types
                    if shape.shape_type == 6:  # group shape
                        continue
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

        master.save(save_path)
        print(f"Combined presentation saved as: {save_path}")

    def Calculate_CA_Average(self, current_df, pulses_to_add=6):
        p1_df = pd.DataFrame()
        p2_df = pd.DataFrame()
        p3_df = pd.DataFrame()
        p4_df = pd.DataFrame()

        p1_sum_df = pd.DataFrame()
        p2_sum_df = pd.DataFrame()
        p3_sum_df = pd.DataFrame()
        p4_sum_df = pd.DataFrame()

        wavelength_array = list(dict.fromkeys([col.split('nm')[0][-3:] for col in current_df.columns]))

        for wavelength in wavelength_array:

            cur_p1 = current_df[f"{wavelength}nm_P1_Power-2600"]
            cur_p2 = current_df[f"{wavelength}nm_P2_Power-2600"]
            cur_p3 = current_df[f"{wavelength}nm_P3_Power-2600"]
            cur_p4 = current_df[f"{wavelength}nm_P4_Power-2600"]

            p1_sums = []
            p2_sums = []
            p3_sums = []
            p4_sums = []

            for i in range(0, pulses_to_add):
                p1_df[f"{wavelength}nm-Pulse_{i}"] = cur_p1.iloc[(0 + i * 90):(50 + i * 90)].reset_index(drop=True)
                p2_df[f"{wavelength}nm-Pulse_{i}"] = cur_p2.iloc[(0 + i * 90):(50 + i * 90)].reset_index(drop=True)
                p3_df[f"{wavelength}nm-Pulse_{i}"] = cur_p3.iloc[(0 + i * 90):(50 + i * 90)].reset_index(drop=True)
                p4_df[f"{wavelength}nm-Pulse_{i}"] = cur_p4.iloc[(0 + i * 90):(50 + i * 90)].reset_index(drop=True)

                p1_sums.extend( [p1_df[f"{wavelength}nm-Pulse_{i}"].iloc[12:48].values.mean() -
                           p1_df[f"{wavelength}nm-Pulse_{i}"].iloc[0:5].values.mean()])
                p2_sums.extend([p2_df[f"{wavelength}nm-Pulse_{i}"].iloc[12:48].mean() -
                           p2_df[f"{wavelength}nm-Pulse_{i}"].iloc[0:5].values.mean()])
                p3_sums.extend([p3_df[f"{wavelength}nm-Pulse_{i}"].iloc[12:48].values.mean() -
                           p3_df[f"{wavelength}nm-Pulse_{i}"].iloc[0:5].values.mean()])
                p4_sums.extend([p4_df[f"{wavelength}nm-Pulse_{i}"].iloc[12:48].values.mean() -
                           p4_df[f"{wavelength}nm-Pulse_{i}"].iloc[0:5].values.mean()])


            p1_sum_df[f"{wavelength}nm-Pulse"] = p1_sums
            p2_sum_df[f"{wavelength}nm-Pulse"] = p2_sums
            p3_sum_df[f"{wavelength}nm-Pulse"] = p3_sums
            p4_sum_df[f"{wavelength}nm-Pulse"] = p4_sums

        return [p1_sum_df, p2_sum_df, p3_sum_df, p4_sum_df], [p1_df, p2_df, p3_df, p4_df]

    def select_folders(self):
        def drop(event):
            raw_paths = event.data
            # Split by whitespace if multiple folders dropped
            paths = root.tk.splitlist(raw_paths)
            folder_paths.extend([p for p in paths if os.path.isdir(p)])
            print("Selected folders:")
            for p in folder_paths:
                print(p)
            root.quit()

        folder_paths = []

        root = TkinterDnD.Tk()
        root.title("Drop folders here or click to select")
        root.geometry("500x200")

        label = tk.Label(root, text="Drop folders here or click this window to select", bg="lightgray")
        label.pack(expand=True, fill="both", padx=20, pady=20)

        # Bind drag-and-drop
        label.drop_target_register(DND_FILES)
        label.dnd_bind('<<Drop>>', drop)

        # Optional: click-to-select fallback
        def open_dialog(event):
            paths = filedialog.askopenfilenames(title="Select files in folders")
            folders = {os.path.dirname(p) for p in paths}
            folder_paths.extend(folders)
            print("Selected folders:")
            for p in folder_paths:
                print(p)
            root.quit()

        label.bind("<Button-1>", open_dialog)

        root.mainloop()

        return list(set(folder_paths))

if __name__ == "__main__":
    # Set the directory where your CSV files are located

    da = Data_Analyzer()

    #da.Generate_Comparison_Slide(da.select_folders())


    for data_folder in da.select_folders():
        da.Generate_Summary_Slide(data_folder)



